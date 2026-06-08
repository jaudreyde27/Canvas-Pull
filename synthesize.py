"""
Synthesize course files into a study guide using the Claude API.

Reads from:
  - downloads/  (output from pull_canvas.py / the pull workflow)
  - source-files/ (manually uploaded files)

Writes:
  - study_guide.md

Required env: ANTHROPIC_API_KEY
Optional env:  COURSE_NAME
"""

import os
import sys
import base64
from pathlib import Path

import anthropic

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None

API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
COURSE_NAME = os.environ.get("COURSE_NAME", "Course")
MAP_MODEL = "claude-haiku-4-5-20251001"      # fast + cheap for per-file extraction
REDUCE_MODEL = "claude-sonnet-4-6"           # high quality for final synthesis

if not API_KEY:
    print("ERROR: ANTHROPIC_API_KEY is not set.")
    sys.exit(1)

client = anthropic.Anthropic(api_key=API_KEY)

# ── Format description (condensed from user's example study guide) ──────────

FORMAT_DESCRIPTION = """
This is a personal reference guide — NOT an exam prep guide. The reader is someone who took the course
and wants a rich, conceptual reminder of what they learned. Write it accordingly: no exam tips, no
"exam note" callouts, no exam framing of any kind.

SECTION STRUCTURE — follow this order strictly for every section:
1. Open by explaining the concept in plain language — what is it, why does it matter, what problem does it solve?
   Use real examples woven into the explanation (not listed in brackets at the end).
2. Once the concept is fully established, introduce any sub-frameworks or nuances.
3. Only after the concept is fully flushed out should you introduce equations or technical specifics.
4. Optionally close with key definitions if they add precision — but definitions should NEVER open a section.

ORGANIZATION:
- Organize sections CHRONOLOGICALLY following the order topics were introduced in the course.
  Use the session numbers in the file names only as a sequencing signal — do NOT reference session
  numbers in the guide itself, and do NOT label sections by session (e.g. no "Session 3-4: Little's Law").
- If a concept spans multiple sessions, consolidate it into one cohesive section with a clean topic title.
- When supplementary files (frameworks, formulas, case studies) relate to a topic, weave them into
  that topic's section rather than appending them at the end.

EXAMPLES — mandatory format:
- Every example must include a one-sentence explanation of HOW the concept applies, not just a name.
- WRONG: "Cost/Performance Leadership — efficiency and scale as primary moat. Example: [Toyota, Walmart]"
- RIGHT: "Cost/Performance Leadership — efficiency and scale as primary moat. Toyota achieves this by
  eliminating waste at every production step so that quality and low cost reinforce each other rather
  than trade off; Walmart does it through distribution scale that lets it undercut rivals on price
  while maintaining margins."

DEPTH AND RICHNESS:
- Each section should paint a full picture: what the concept is, why it exists, how it works in practice,
  what the key tensions or tradeoffs are, and where it breaks down or gets complicated.
- Do not reduce concepts to a bullet list. Use prose where it earns its place, especially when explaining
  WHY something works the way it does.
- Capture the intellectual richness of the source material — the mechanisms, the edge cases, the
  counterintuitive insights.

STYLE RULES:
- Bold every key term on first use
- Arrow notation (→) for causal relationships
- "Key distinction:" prefix for clarifications between commonly confused concepts
- Number all top-level sections sequentially (1, 2, 3 …)
- NO exam notes, exam tips, exam traps, or any exam framing whatsoever
"""

FORMAT_EXAMPLE = """
EXAMPLE OF THE TARGET STYLE (from a different course — match this depth, flow, and richness):

## 3. Competitive Positioning & Moats

The central challenge for any firm is not just creating value — it is keeping it. A firm that creates
enormous value but operates in a market with perfect competition will price down to cost and capture
nothing for itself. **Positioning advantage** is the mechanism by which firms reduce the price pressure
they face, not by negotiating harder, but by finding a market position where direct comparison to
rivals becomes harder for buyers to make.

The intuition is geographic before it is strategic: two coffee shops on opposite sides of a city do not
compete with each other the way two shops on the same block do. Walmart understood this when it expanded
into rural towns too small to support two discount retailers — by being first into markets where a second
entrant was economically unviable, it locked in a position before competition could arrive.

**Product differentiation** works on the same principle but along taste rather than geography. When
AbInBev tries to compete with craft beer, it runs into a perception problem: even if it brews a
technically similar product, customers who value authenticity and local identity see it as a fundamentally
different (and inferior) substitute. The differentiation is as much about story and identity as
about the liquid in the glass.

### What Protects a Position — Moats

A positioning advantage is only durable if something prevents rivals from copying it. High returns
attract entry; without a barrier, the advantage erodes. The most important moats are:

**Economies of scale** operate when fixed costs are large relative to variable costs, so that the firm
already at scale can produce at a cost that makes entry irrational. The key metric is the ratio of
Minimum Efficient Scale to total market size — when one firm's MES is close to the whole market,
a second firm simply cannot reach the same cost structure. This is why small cities support only one
bike-share operator while large cities can sustain several.

**Network effects** create a self-reinforcing dynamic where each new user increases the value of
the product for every existing user. An entrant without an installed base faces a value creation
disadvantage before the first price negotiation even begins.

Key distinction: A product going viral is NOT the same as a product having network effects. Poppi
beverage spread through social media — but more people buying Poppi does not make Poppi better for
any individual buyer. TikTok, by contrast, becomes genuinely more valuable as more creators join,
because each new creator adds content that increases time-on-platform for everyone else.
"""


# ── File extraction helpers ──────────────────────────────────────────────────

def extract_pdf_text(path: Path) -> str:
    if fitz is None:
        return f"[PDF could not be read — PyMuPDF unavailable]"
    doc = fitz.open(str(path))
    pages = [page.get_text() for page in doc]
    text = "\n".join(pages)
    return text[:40000]  # cap per file


def extract_pptx_text(path: Path) -> str:
    if PptxPresentation is None:
        return "[PPTX could not be read — python-pptx unavailable]"
    prs = PptxPresentation(str(path))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)[:40000]


def encode_image(path: Path) -> tuple[str, str]:
    media = {".jpg": "image/jpeg", ".jpeg": "image/jpeg",
             ".png": "image/png", ".gif": "image/gif", ".webp": "image/webp"}
    mt = media.get(path.suffix.lower(), "image/jpeg")
    data = base64.standard_b64encode(path.read_bytes()).decode()
    return data, mt


# ── Map phase: extract key points from one file ──────────────────────────────

def extract_key_points(path: Path) -> str:
    name = path.name
    suffix = path.suffix.lower()
    image_exts = {".png", ".jpg", ".jpeg", ".gif", ".webp"}

    print(f"  → {name}")

    extraction_instructions = (
        "Extract comprehensive notes from this course material. For each concept capture:\n"
        "- The core idea in plain language — what it is and why it matters\n"
        "- The mechanisms and logic behind it (not just what, but why and how)\n"
        "- Any formulas or equations, clearly labeled, with what each variable means\n"
        "- Real-world examples including a brief explanation of how the concept applies\n"
        "- Key tensions, tradeoffs, or counterintuitive insights\n"
        "- How this topic connects to other topics in the material\n"
        "Preserve all terminology exactly. Be thorough — capture the richness, not just the surface.\n"
        "Note the session number if visible in the file name or content, as it will be used for ordering.\n"
    )

    if suffix == ".pdf":
        text = extract_pdf_text(path)
        content = [{"type": "text", "text": (
            f"Course file: '{name}'\n\n{extraction_instructions}\n\n{text}"
        )}]

    elif suffix in {".pptx", ".ppt"}:
        text = extract_pptx_text(path)
        content = [{"type": "text", "text": (
            f"Course slide deck: '{name}'\n\n{extraction_instructions}\n\n{text}"
        )}]

    elif suffix in image_exts:
        img_data, media_type = encode_image(path)
        content = [
            {"type": "image", "source": {"type": "base64", "media_type": media_type, "data": img_data}},
            {"type": "text", "text": (
                f"This is a course slide/image: '{name}'. {extraction_instructions}"
            )},
        ]
    else:
        return ""

    resp = client.messages.create(
        model=MAP_MODEL,
        max_tokens=2048,
        messages=[{"role": "user", "content": content}],
    )
    return resp.content[0].text


# ── Reduce phase: synthesize all notes into final study guide ────────────────

def synthesize(all_notes: list[tuple[str, str]]) -> str:
    print("\nSynthesizing into study guide …")

    combined = "\n\n".join(
        f"=== SOURCE: {fname} ===\n{notes}"
        for fname, notes in all_notes
    )

    # Trim if very large (> 150k chars ≈ 37k tokens)
    if len(combined) > 150000:
        combined = combined[:150000] + "\n\n[... additional material truncated for context limits ...]"

    prompt = (
        f'You are creating a comprehensive personal reference guide for "{COURSE_NAME}".\n\n'
        f"{FORMAT_DESCRIPTION}\n\n"
        f"{FORMAT_EXAMPLE}\n\n"
        "Now synthesize the course notes below into a complete study guide following those guidelines exactly.\n\n"
        "Critical requirements:\n"
        "- CHRONOLOGICAL ORDER: Use the session numbers in the source file names (Session 1, Session 2, etc.)\n"
        "  as your primary ordering signal. Supplementary files (formulas, frameworks, case studies) should\n"
        "  be woven into whichever session section they belong to — not appended at the end.\n"
        "- CONCEPT BEFORE TECHNICAL: For every topic, fully explain what it is and why it matters in plain\n"
        "  language first. Only introduce equations, formulas, or calculations after the concept is established.\n"
        "- RICH EXAMPLES: Every example must include a one-sentence explanation of how the concept applies\n"
        "  to that specific company or situation — never just a name in brackets.\n"
        "- DEPTH: Do not flatten concepts into thin bullet lists. Capture the mechanisms, the tensions,\n"
        "  the counterintuitive insights, and the edge cases. This guide should feel intellectually rich.\n"
        "- NO EXAM FRAMING: No exam notes, exam tips, exam traps, or any exam language whatsoever.\n"
        "- Do not repeat content; consolidate when the same idea appears in multiple files.\n\n"
        "COURSE NOTES:\n\n"
        f"{combined}"
    )

    resp = client.messages.create(
        model=REDUCE_MODEL,
        max_tokens=8192,
        messages=[{"role": "user", "content": prompt}],
    )
    return resp.content[0].text


# ── Main ─────────────────────────────────────────────────────────────────────

SUPPORTED = {".pdf", ".pptx", ".ppt", ".png", ".jpg", ".jpeg", ".gif", ".webp"}


def collect_files(*dirs: Path) -> list[Path]:
    files = []
    for d in dirs:
        if d.exists():
            for f in sorted(d.rglob("*")):
                if f.is_file() and f.suffix.lower() in SUPPORTED:
                    files.append(f)
    return files


def main() -> None:
    files = collect_files(Path("downloads"), Path("source-files"))

    if not files:
        print("No supported files found in downloads/ or source-files/")
        print("Run the 'Pull Canvas Files' workflow first, or add files to source-files/")
        sys.exit(1)

    print(f"Found {len(files)} file(s)\n")

    all_notes: list[tuple[str, str]] = []
    for f in files:
        try:
            notes = extract_key_points(f)
            if notes.strip():
                all_notes.append((f.name, notes))
        except Exception as exc:
            print(f"  [error] {f.name}: {exc}")

    if not all_notes:
        print("No content could be extracted from any file.")
        sys.exit(1)

    guide = synthesize(all_notes)

    out = Path("study_guide.md")
    out.write_text(guide, encoding="utf-8")
    size = len(guide)
    print(f"\nDone — study_guide.md ({size:,} chars)")


if __name__ == "__main__":
    main()
